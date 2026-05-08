"""Generate the Testcontainers demo presentation using the Red Hat template."""

import copy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

TEMPLATE = Path("/home/rauherna/Downloads/Red Hat standard presentation template - Please make a copy.pptx")
OUTPUT = Path(__file__).parent / "Testcontainers - Integration Testing Demo.pptx"
ARCH_IMAGE = Path(__file__).parent / "testcontainers-architecture.png"
MOCK_IMAGE = Path(__file__).parent / "mock-call-flow.png"

RED_HAT_RED = RGBColor(0xEE, 0x00, 0x00)
RED_HAT_BLACK = RGBColor(0x15, 0x15, 0x15)
RED_HAT_GRAY = RGBColor(0x6A, 0x6E, 0x73)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG = RGBColor(0xF5, 0xF5, 0xF5)


def clear_slides(prs):
    """Remove all existing slides from the presentation."""
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=RED_HAT_BLACK, alignment=PP_ALIGN.LEFT,
                 font_name="Red Hat Display"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_code_box(slide, left, top, width, height, code_text, font_size=11):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = code_text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = "Red Hat Mono"
        p.font.color.rgb = RED_HAT_BLACK
        p.space_after = Pt(2)

    # Add background fill to the textbox
    fill = txBox.fill
    fill.solid()
    fill.fore_color.rgb = CODE_BG
    return txBox


def add_bullet_slide(prs, title_text, bullets, section_marker=""):
    """Add a slide with title and bullet points."""
    slide_layout = prs.slide_layouts[11]  # CUSTOM_1_1_1_1_1_1_1_1 - simple title+subtitle
    slide = prs.slides.add_slide(slide_layout)

    # Clear all placeholders
    for ph in slide.placeholders:
        ph.text = ""

    # Title
    add_text_box(slide, Cm(1.5), Cm(1.0), Cm(20), Cm(1.5),
                 title_text, font_size=28, bold=True, color=RED_HAT_BLACK)

    # Section marker (small red text top-right)
    if section_marker:
        add_text_box(slide, Cm(1.5), Cm(0.3), Cm(10), Cm(0.8),
                     section_marker, font_size=10, color=RED_HAT_RED)

    # Bullets
    bullet_top = Cm(3.5)
    for bullet in bullets:
        # Red triangle marker
        add_text_box(slide, Cm(1.5), bullet_top, Cm(0.8), Cm(0.7),
                     "\u25B6", font_size=10, color=RED_HAT_RED)
        add_text_box(slide, Cm(2.5), bullet_top, Cm(20), Cm(1.2),
                     bullet, font_size=16, color=RED_HAT_BLACK)
        bullet_top += Cm(1.8)

    return slide


def create_title_slide(prs):
    """Slide 1: Title slide."""
    slide_layout = prs.slide_layouts[3]  # TITLE_1_2 - title with subtitle
    slide = prs.slides.add_slide(slide_layout)

    for ph in slide.placeholders:
        ph.text = ""

    add_text_box(slide, Cm(1.5), Cm(5.0), Cm(22), Cm(3.0),
                 "Testcontainers", font_size=44, bold=True, color=RED_HAT_BLACK)
    add_text_box(slide, Cm(1.5), Cm(8.5), Cm(22), Cm(1.5),
                 "Integration Testing with Real Dependencies",
                 font_size=22, color=RED_HAT_GRAY)
    add_text_box(slide, Cm(1.5), Cm(11.0), Cm(15), Cm(1.5),
                 "Raúl Hernández\nO-Cloud Manager Team",
                 font_size=14, color=RED_HAT_GRAY)


def create_agenda_slide(prs):
    """Slide 2: Agenda."""
    bullets = [
        "The problem: integration testing with external dependencies",
        "Testcontainers: what it is and how it works",
        "Architecture: real-world example from our codebase",
        "Appendix: code walkthrough (3 steps)",
    ]
    add_bullet_slide(prs, "Agenda", bullets)


def create_problem_slide(prs):
    """Slide 3: The problem."""
    bullets = [
        "Unit tests mock everything — miss real integration bugs",
        "Shared test databases create flaky, order-dependent tests",
        "CI environments differ from local — \"works on my machine\"",
        "Manual Docker setup for tests is brittle and not portable",
    ]
    add_bullet_slide(prs, "The problem", bullets, section_marker="Context")


def create_solution_slide(prs):
    """Slide 4: What is testcontainers."""
    bullets = [
        "Library that manages Docker containers in test code",
        "Each test gets a fresh, isolated container (DB, broker, etc.)",
        "Containers start before test, terminate after — automatic lifecycle",
        "Modules for popular services: PostgreSQL, Redis, Kafka, ...",
        "Available in Go, Java, Python, .NET, Rust, Node.js",
    ]
    add_bullet_slide(prs, "Testcontainers", bullets, section_marker="Solution")


def create_mock_slide(prs):
    """Slide: How mocking works (simple call flow)."""
    slide_layout = prs.slide_layouts[11]
    slide = prs.slides.add_slide(slide_layout)

    for ph in slide.placeholders:
        ph.text = ""

    add_text_box(slide, Cm(1.5), Cm(0.3), Cm(10), Cm(0.8),
                 "Context", font_size=10, color=RED_HAT_RED)
    add_text_box(slide, Cm(1.5), Cm(1.0), Cm(22), Cm(1.5),
                 "How mocking works today", font_size=28, bold=True,
                 color=RED_HAT_BLACK)

    if MOCK_IMAGE.exists():
        img_left = Cm(5.0)
        img_top = Cm(4.0)
        img_width = Cm(15.0)
        slide.shapes.add_picture(str(MOCK_IMAGE), img_left, img_top, width=img_width)

    add_text_box(slide, Cm(1.5), Cm(13.0), Cm(22), Cm(1.5),
                 "The mock returns canned data — no real SQL, no schema validation, no constraint checks.",
                 font_size=13, color=RED_HAT_GRAY)


def create_architecture_slide(prs):
    """Slide 5: Architecture diagram."""
    slide_layout = prs.slide_layouts[11]
    slide = prs.slides.add_slide(slide_layout)

    for ph in slide.placeholders:
        ph.text = ""

    add_text_box(slide, Cm(1.5), Cm(0.3), Cm(10), Cm(0.8),
                 "Architecture", font_size=10, color=RED_HAT_RED)
    add_text_box(slide, Cm(1.5), Cm(1.0), Cm(22), Cm(1.5),
                 "End-to-end test architecture", font_size=28, bold=True,
                 color=RED_HAT_BLACK)

    # Insert the architecture diagram image
    if ARCH_IMAGE.exists():
        img_left = Cm(2.0)
        img_top = Cm(3.5)
        img_width = Cm(21.0)
        slide.shapes.add_picture(str(ARCH_IMAGE), img_left, img_top, width=img_width)


def create_step1_slide(prs):
    """Slide 6: Appendix - Step 1."""
    slide_layout = prs.slide_layouts[11]
    slide = prs.slides.add_slide(slide_layout)

    for ph in slide.placeholders:
        ph.text = ""

    add_text_box(slide, Cm(1.5), Cm(0.3), Cm(10), Cm(0.8),
                 "Appendix", font_size=10, color=RED_HAT_RED)
    add_text_box(slide, Cm(1.5), Cm(1.0), Cm(22), Cm(1.5),
                 "Step 1: Start a PostgreSQL container",
                 font_size=24, bold=True, color=RED_HAT_BLACK)

    code = '''\
pc, err := postgres.Run(ctx,
    "docker.io/postgres:16-alpine",
    postgres.WithDatabase("resources_test"),
    postgres.WithUsername("test"),
    postgres.WithPassword("test"),
    testcontainers.WithWaitStrategy(
        wait.ForLog("database system is ready").
            WithOccurrence(2).
            WithStartupTimeout(30*time.Second),
    ),
)'''
    add_code_box(slide, Cm(1.0), Cm(3.0), Cm(22), Cm(8.0), code, font_size=12)

    add_text_box(slide, Cm(1.0), Cm(11.5), Cm(22), Cm(2.0),
                 "postgres is a testcontainers module — pre-configured for PostgreSQL.\n"
                 "WaitStrategy ensures the container is fully ready before tests proceed.",
                 font_size=12, color=RED_HAT_GRAY)


def create_step2_slide(prs):
    """Slide 7: Appendix - Step 2."""
    slide_layout = prs.slide_layouts[11]
    slide = prs.slides.add_slide(slide_layout)

    for ph in slide.placeholders:
        ph.text = ""

    add_text_box(slide, Cm(1.5), Cm(0.3), Cm(10), Cm(0.8),
                 "Appendix", font_size=10, color=RED_HAT_RED)
    add_text_box(slide, Cm(1.5), Cm(1.0), Cm(22), Cm(1.5),
                 "Step 2: Run database migrations",
                 font_size=24, bold=True, color=RED_HAT_BLACK)

    code = '''\
// Get connection string from the running container
connStr, _ := pc.ConnectionString(ctx, "sslmode=disable")

// Create connection pool — same as production code
pool, _ := pgxpool.New(ctx, connStr)

// Apply schema migrations from embedded SQL files
source, _ := iofs.New(migrationsFS, "db/migrations")
m, _ := migrate.NewWithSourceInstance("iofs", source,
    migrateConnStr(connStr))
m.Up()'''
    add_code_box(slide, Cm(1.0), Cm(3.0), Cm(22), Cm(7.0), code, font_size=12)

    add_text_box(slide, Cm(1.0), Cm(11.0), Cm(22), Cm(2.0),
                 "Same migration code used in production — guarantees schema parity.\n"
                 "Each test run starts with a clean, correctly-migrated database.",
                 font_size=12, color=RED_HAT_GRAY)


def create_step3_slide(prs):
    """Slide 8: Appendix - Step 3."""
    slide_layout = prs.slide_layouts[11]
    slide = prs.slides.add_slide(slide_layout)

    for ph in slide.placeholders:
        ph.text = ""

    add_text_box(slide, Cm(1.5), Cm(0.3), Cm(10), Cm(0.8),
                 "Appendix", font_size=10, color=RED_HAT_RED)
    add_text_box(slide, Cm(1.5), Cm(1.0), Cm(22), Cm(1.5),
                 "Step 3: Test and cleanup",
                 font_size=24, bold=True, color=RED_HAT_BLACK)

    code = '''\
var _ = AfterSuite(func() {
    cancel()          // signal collector to stop
    testServer.Close() // stop HTTP server
    testEnv.Stop()    // stop envtest K8s API

    // One call tears down the container completely
    Expect(pgContainer.Terminate(ctx)).To(Succeed())
})'''
    add_code_box(slide, Cm(1.0), Cm(3.0), Cm(22), Cm(5.5), code, font_size=12)

    add_text_box(slide, Cm(1.0), Cm(9.0), Cm(22), Cm(2.5),
                 "Terminate() stops the container and releases all resources.\n"
                 "No leftover state between test runs — fully hermetic.",
                 font_size=12, color=RED_HAT_GRAY)

    # Key takeaways
    add_text_box(slide, Cm(1.0), Cm(12.0), Cm(22), Cm(1.5),
                 "Key takeaway: real dependencies, zero maintenance, total isolation.",
                 font_size=14, bold=True, color=RED_HAT_RED)


def create_links_slide(prs):
    """Slide 9: Resources/links."""
    bullets = [
        "testcontainers.com — official site & docs",
        "github.com/testcontainers/testcontainers-go — Go module",
        "Branch: demo.testcontainers — working example in our repo",
    ]
    add_bullet_slide(prs, "Resources", bullets)


def main():
    prs = Presentation(str(TEMPLATE))
    clear_slides(prs)

    create_title_slide(prs)
    create_agenda_slide(prs)
    create_problem_slide(prs)
    create_mock_slide(prs)
    create_solution_slide(prs)
    create_architecture_slide(prs)
    create_step1_slide(prs)
    create_step2_slide(prs)
    create_step3_slide(prs)
    create_links_slide(prs)

    prs.save(str(OUTPUT))
    print(f"Presentation saved to: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
